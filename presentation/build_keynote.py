#!/usr/bin/env python3
"""Generate the three-minute MagPilot assistive-robotics keynote.

The deck uses eight story beats across ten remote-controlled frames, 113
seconds of spoken material, and 67 seconds of auto-playing demo GIFs. The
complete run is exactly three minutes.

Regenerate with:
    python3 presentation/build_keynote.py
"""

import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt


HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(HERE, "assets")
DOCS = os.path.join(ROOT, "docs")

INK = RGBColor(0x10, 0x2A, 0x43)
NAVY = RGBColor(0x09, 0x1A, 0x2B)
BLUE_DK = RGBColor(0x06, 0x3A, 0x6E)
BLUE = RGBColor(0x0A, 0x84, 0xFF)
SKY = RGBColor(0xEA, 0xF3, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xB9, 0xD3, 0xE8)
SUBTLE = RGBColor(0x5C, 0x74, 0x89)
GREEN = RGBColor(0x2E, 0xC4, 0x6B)
WARM = RGBColor(0xFF, 0xC8, 0x57)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def make_slide(bg=SKY):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def add_notes(s, body, cue=None):
    """Add a concise, readable Presenter View script."""
    lines = [line.strip() for line in body.strip().splitlines() if line.strip()]
    entries = [lines[0], " ".join(lines[1:])]
    if cue:
        entries.append(cue)
    tf = s.notes_slide.notes_text_frame
    tf.clear()
    for index, line in enumerate(entries):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(8 if index else 12)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(14 if index != 1 else 18)
        run.font.bold = index != 1
        run.font.color.rgb = BLUE if index == 2 else INK


def box(s, left, top, width, height, color, radius=False, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shape_type, left, top, width, height)
    if radius:
        shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def set_opacity(shp, opacity):
    solid = shp.fill._xPr.find(qn("a:solidFill"))
    color = solid.find(qn("a:srgbClr"))
    color.append(color.makeelement(qn("a:alpha"), {
        "val": str(int(max(0.0, min(1.0, opacity)) * 100000))
    }))


def add_text(s, left, top, width, height, paragraphs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=4, line_spacing=1.0):
    """Add styled text.

    ``paragraphs`` is a list of paragraphs. Each paragraph is either one
    ``(text, size, bold, color)`` tuple or a list of those tuples.
    """
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    for index, paragraph in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        runs = [paragraph] if isinstance(paragraph, tuple) else paragraph
        for value, size, bold, color in runs:
            run = p.add_run()
            run.text = value
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return tb


def fit_image(path, left, top, width, height):
    iw, ih = Image.open(path).size
    image_ratio = iw / ih
    box_ratio = width / height
    if image_ratio > box_ratio:
        w = width
        h = int(width / image_ratio)
    else:
        h = height
        w = int(height * image_ratio)
    return (
        int(left + (width - w) / 2),
        int(top + (height - h) / 2),
        int(w),
        int(h),
    )


def add_picture(s, path, left, top, width, height):
    l, t, w, h = fit_image(path, left, top, width, height)
    return s.shapes.add_picture(path, l, t, w, h)


def add_cover(s, path, left, top, width, height):
    iw, ih = Image.open(path).size
    image_ratio = iw / ih
    box_ratio = width / height
    pic = s.shapes.add_picture(path, left, top, width, height)
    if image_ratio > box_ratio:
        crop = (1 - box_ratio / image_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - image_ratio / box_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def kicker(s, value, dark=False, left=Inches(0.78), top=Inches(0.55)):
    add_text(
        s, left, top, Inches(8.5), Inches(0.35),
        [(value.upper(), 13, True, WHITE if dark else BLUE)],
    )


def footer(s, number, dark=False):
    color = MIST if dark else SUBTLE
    add_text(
        s, Inches(0.6), Inches(7.08), Inches(4), Inches(0.25),
        [("MagPilot | TUM ARIES Lab", 9, False, color)],
    )
    add_text(
        s, Inches(12), Inches(7.08), Inches(0.7), Inches(0.25),
        [(str(number), 9, True, color)],
        align=PP_ALIGN.RIGHT,
    )


def add_click_transition(s, effect="fade", direction=None):
    """Add a click-advanced PowerPoint transition to this slide."""
    attrs = ' spd="fast" advClick="1"'
    child_attrs = f' dir="{direction}"' if direction else ""
    transition = parse_xml(
        f'<p:transition {nsdecls("p")}{attrs}>'
        f'<p:{effect}{child_attrs}/>'
        '</p:transition>'
    )
    slide_element = s._element
    for index, child in enumerate(slide_element):
        if child.tag in (qn("p:timing"), qn("p:extLst")):
            slide_element.insert(index, transition)
            break
    else:
        slide_element.append(transition)


def metric(s, left, top, number, label, dark=False):
    primary = WHITE if dark else INK
    secondary = MIST if dark else SUBTLE
    add_text(
        s, left, top, Inches(3.5), Inches(0.72),
        [(number, 28, True, BLUE)],
    )
    add_text(
        s, left, top + Inches(0.65), Inches(3.5), Inches(0.75),
        [(label, 14, False, secondary)],
        line_spacing=1.05,
    )
    return primary


SETUP = os.path.join(DOCS, "setup.png")
WRITING_UI = os.path.join(DOCS, "interface.png")
UI = os.path.join(DOCS, "magpilot.png")
LAUNCHER = os.path.join(DOCS, "launcher.png")
ACTION_MAPPING = os.path.join(DOCS, "action_mapping.png")
LOGO = os.path.join(ASSETS, "logo.png")
TUM_LOGO = os.path.join(ASSETS, "tum-logo.png")
MIRMI_LOGO = os.path.join(DOCS, "affiliations", "mirmi-logo.png")
ASSISTIVE_CONCEPT = os.path.join(ASSETS, "assistive_concept.png")
GIF_CLASSIFY = os.path.join(ASSETS, "demo_classify_slide.gif")
GIF_TELEOP = os.path.join(ASSETS, "demo_teleop_slide.gif")


# 1. Title and promise
s = make_slide(NAVY)
add_cover(s, SETUP, 0, 0, SW, SH)
scrim = box(s, 0, 0, Inches(7.45), SH, NAVY)
set_opacity(scrim, 0.89)
add_picture(
    s, LOGO, Inches(0.8), Inches(0.62), Inches(0.82), Inches(0.82)
)
add_text(
    s, Inches(0.82), Inches(1.58), Inches(5.9), Inches(0.35),
    [("INTRODUCING", 13, True, BLUE)],
)
add_text(
    s, Inches(0.8), Inches(2.02), Inches(6.25), Inches(1.2),
    [("MagPilot", 72, True, WHITE)],
)
add_text(
    s, Inches(0.82), Inches(3.42), Inches(6.35), Inches(1.35),
    [
        ("Robot control from", 31, False, MIST),
        ("one tiny passive magnet.", 31, True, WHITE),
    ],
    space_after=4,
)
add_text(
    s, Inches(0.82), Inches(6.76), Inches(6.3), Inches(0.35),
    [("COLMAG | TUM ARIES Lab | Simon Kruelle", 11, False, MIST)],
)
add_notes(s, """
[0:00-0:08]
Introducing MagPilot: robot control from one tiny passive magnet. Hold up the
magnet and pause for a moment.
""", cue="REMOTE: click to introduce the problem.")


# 2. Problem and assistive concept, built in two remote clicks
def add_concept_base():
    slide = make_slide(NAVY)
    add_cover(slide, ASSISTIVE_CONCEPT, 0, 0, SW, SH)
    concept_scrim = box(slide, 0, 0, Inches(4.18), SH, NAVY)
    set_opacity(concept_scrim, 0.93)
    add_text(
        slide, Inches(0.58), Inches(6.82), Inches(3.1), Inches(0.36),
        [("AI-GENERATED CONCEPT | NOT USER-VALIDATED", 7.5, True, MIST)],
    )
    add_text(
        slide, Inches(12.0), Inches(7.08), Inches(0.7), Inches(0.25),
        [("2", 9, True, MIST)],
        align=PP_ALIGN.RIGHT,
    )
    return slide


s = add_concept_base()
kicker(
    s, "Why MagPilot", dark=True,
    left=Inches(0.58), top=Inches(0.48),
)
add_text(
    s, Inches(0.58), Inches(1.2), Inches(3.15), Inches(1.55),
    [
        ("Robot arms can assist.", 29, False, WHITE),
        ("Controllers can exclude.", 29, True, BLUE),
    ],
    space_after=2,
)
box(s, Inches(0.6), Inches(3.35), Inches(0.1), Inches(1.18), BLUE)
add_text(
    s, Inches(0.92), Inches(3.34), Inches(2.9), Inches(1.2),
    [
        ("CAN ONE TINY MAGNET", 11, True, MIST),
        ("BECOME THE WHOLE INTERFACE?", 18, True, WHITE),
    ],
    space_after=7,
)
add_click_transition(s, effect="fade")
add_notes(s, """
[0:08-0:20]
Robot arms can give people more physical independence, but their controllers
still assume capable hands. MagPilot asks: can one tiny passive magnet become
the whole interface?
""", cue="REMOTE: click to reveal magnet placement and benefits.")


s = add_concept_base()
kicker(
    s, "One tiny neodymium magnet", dark=True,
    left=Inches(0.58), top=Inches(0.48),
)
add_text(
    s, Inches(0.58), Inches(1.08), Inches(3.15), Inches(1.0),
    [("Place it almost anywhere.", 28, True, WHITE)],
)
add_text(
    s, Inches(0.6), Inches(2.18), Inches(3.05), Inches(0.82),
    [
        ("BANDAGE  |  GLOVE", 12, True, BLUE),
        ("CUFF  |  RESIDUAL LIMB", 12, True, BLUE),
    ],
    space_after=5,
)
add_text(
    s, Inches(0.6), Inches(3.22), Inches(3.05), Inches(0.68),
    [
        ("LOCATION  >  MOTION", 13, True, WHITE),
        ("TILT  >  GRIP", 13, True, WHITE),
    ],
    space_after=5,
)
box(s, Inches(0.6), Inches(4.2), Inches(3.05), Inches(0.08), BLUE)
add_text(
    s, Inches(0.6), Inches(4.52), Inches(3.15), Inches(0.82),
    [
        ("NO BATTERY  |  NO CABLES", 12, True, MIST),
        ("NO BUTTONS  |  NO GRIP", 12, True, MIST),
    ],
    space_after=6,
)
add_text(
    s, Inches(0.6), Inches(5.66), Inches(3.1), Inches(0.48),
    [("PURE MAGNETIC MAGIC.", 16, True, BLUE)],
)
add_click_transition(s, effect="wipe", direction="u")
add_notes(s, """
[0:20-0:38]
A small neodymium magnet can sit almost anywhere: on a bandage, glove, cuff, or
residual limb. We only need its location, plus tilt for gripping. No battery,
no cables, no buttons, no grip. Pure magnetic magic.
""", cue="REMOTE: click to highlight classified commands.")


# 3. Two operating modes, with the emphasis moved by the remote
def add_modes_frame(active_mode):
    slide = make_slide(SKY)
    kicker(slide, "One magnet. Two modes.")
    add_text(
        slide, Inches(0.78), Inches(1.0), Inches(11.8), Inches(0.58),
        [("Choose a task or pilot continuously.", 29, True, INK)],
    )
    modes = [
        (
            "commands", Inches(0.48), "01  CLASSIFIED COMMANDS",
            WRITING_UI, "WRITE  >  RECOGNIZE  >  RUN",
        ),
        (
            "teleop", Inches(6.77), "02  MAGPILOT TELEOPERATION",
            UI, "MOVE  >  TILT  >  GRIP  >  PAUSE",
        ),
    ]
    for key, left, label, image_path, flow in modes:
        is_active = key == active_mode
        card = box(
            slide, left, Inches(1.78), Inches(6.08), Inches(5.02),
            WHITE, radius=True, line=BLUE if is_active else None,
        )
        if is_active:
            card.line.width = Pt(2.5)
            box(
                slide, left + Inches(0.18), Inches(1.78),
                Inches(5.72), Inches(0.08), BLUE,
            )
        add_text(
            slide, left + Inches(0.24), Inches(2.08),
            Inches(5.6), Inches(0.34),
            [(label, 13, True, BLUE if is_active else SUBTLE)],
        )
        add_picture(
            slide, image_path, left + Inches(0.2), Inches(2.58),
            Inches(5.68), Inches(3.47),
        )
        if not is_active:
            shade = box(
                slide, left + Inches(0.2), Inches(2.58),
                Inches(5.68), Inches(3.47), SKY,
            )
            set_opacity(shade, 0.72)
        add_text(
            slide, left + Inches(0.24), Inches(6.28),
            Inches(5.6), Inches(0.3),
            [(flow, 11, True, BLUE if is_active else SUBTLE)],
            align=PP_ALIGN.CENTER,
        )
    footer(slide, 3)
    return slide


s = add_modes_frame("commands")
add_click_transition(s, effect="fade")
add_notes(s, """
[0:38-0:50]
First, command mode. Write a letter or digit; the classifier recognizes it and
runs the action mapped to that symbol.
""", cue="REMOTE: click to move the highlight to Teleoperation.")


s = add_modes_frame("teleop")
add_click_transition(s, effect="fade")
add_notes(s, """
[0:50-1:02]
Second, MagPilot Teleoperation. Position and height follow the magnet, tilt
controls the gripper, and lifting away pauses the arm. The magnet itself
switches between both modes.
""", cue="REMOTE: click once to start the first demonstration.")


# 4. Auto-playing classified-command demo
s = make_slide(NAVY)
s.shapes.add_picture(GIF_CLASSIFY, 0, 0, SW, SH)
top_band = box(s, 0, 0, SW, Inches(1.03), NAVY)
set_opacity(top_band, 0.92)
add_text(
    s, Inches(0.68), Inches(0.18), Inches(9.8), Inches(0.7),
    [
        ("PROOF 1  ", 13, True, BLUE),
        ("One movement becomes a command.", 26, True, WHITE),
    ],
    anchor=MSO_ANCHOR.MIDDLE,
)
badge = box(
    s, Inches(10.62), Inches(0.23), Inches(1.98), Inches(0.52), BLUE, radius=True
)
add_text(
    s, Inches(10.62), Inches(0.23), Inches(1.98), Inches(0.52),
    [("AUTO PLAY  0:39", 10, True, WHITE)],
    align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE,
)
add_notes(s, """
[1:02-1:41]
Do not speak. Let the 39-second classified-command demonstration play.
""", cue="REMOTE: when the clip finishes, click once.")


# 5. Auto-playing MagPilot demo
s = make_slide(NAVY)
s.shapes.add_picture(GIF_TELEOP, 0, 0, SW, SH)
top_band = box(s, 0, 0, SW, Inches(1.03), NAVY)
set_opacity(top_band, 0.92)
add_text(
    s, Inches(0.68), Inches(0.18), Inches(9.8), Inches(0.7),
    [
        ("PROOF 2  ", 13, True, BLUE),
        ("MagPilot Teleoperation picks up an object.", 26, True, WHITE),
    ],
    anchor=MSO_ANCHOR.MIDDLE,
)
badge = box(
    s, Inches(10.62), Inches(0.23), Inches(1.98), Inches(0.52), BLUE, radius=True
)
add_text(
    s, Inches(10.62), Inches(0.23), Inches(1.98), Inches(0.52),
    [("AUTO PLAY  0:28", 10, True, WHITE)],
    align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE,
)
add_notes(s, """
[1:41-2:09]
Do not speak. Let the 28-second MagPilot Teleoperation demonstration play.
""", cue="REMOTE: when the clip finishes, click once.")


# 6-7. Customer journey with a click-revealed customization screen
def add_platform_base(s, slide_number):
    kicker(
        s, "From download to deployment",
        left=Inches(0.68), top=Inches(0.55),
    )
    add_text(
        s, Inches(0.68), Inches(1.02), Inches(4.18), Inches(1.35),
        [
            ("Software first.", 30, False, INK),
            ("Hardware when it fits.", 30, True, INK),
        ],
        space_after=3,
    )
    layers = [
        ("1", "TRY NOW", "Trackpad + Gazebo"),
        ("2", "ADD SENSING", "Our board or tailored hardware"),
        ("3", "DEPLOY", "FR3 | Panda | commercial robots"),
    ]
    for index, (number, title, detail) in enumerate(layers):
        top = Inches(2.55 + index * 1.1)
        add_text(
            s, Inches(0.72), top, Inches(0.48), Inches(0.48),
            [(number, 23, True, BLUE)],
        )
        add_text(
            s, Inches(1.3), top + Inches(0.02), Inches(3.2), Inches(0.3),
            [(title, 13, True, INK)],
        )
        add_text(
            s, Inches(1.3), top + Inches(0.38), Inches(3.25), Inches(0.48),
            [(detail, 13, False, SUBTLE)],
        )
    box(
        s, Inches(4.92), Inches(0.62), Inches(7.8), Inches(6.17),
        WHITE, radius=True,
    )
    add_picture(
        s, LAUNCHER, Inches(5.1), Inches(0.82), Inches(7.44), Inches(5.5)
    )
    add_text(
        s, Inches(5.16), Inches(6.42), Inches(7.3), Inches(0.27),
        [("ONE SOFTWARE STACK  >  ANY SUPPORTED ROBOT", 10, True, BLUE)],
        align=PP_ALIGN.CENTER,
    )
    footer(s, slide_number)


s = make_slide(SKY)
add_platform_base(s, 6)
add_notes(s, """
[2:09-2:29]
The software can ship today. Customers first try MagPilot with a trackpad and
Gazebo, without hardware. If the workflow fits, they add our sensor board or a
tailored setup. The same stack then moves to an FR3, Panda, or another
commercial robot we integrate.
""", cue="REMOTE: click to drop in the customizable action screen.")


s = make_slide(SKY)
add_platform_base(s, 7)
window = box(
    s, Inches(6.35), Inches(1.15), Inches(5.82), Inches(5.08),
    NAVY, radius=True,
)
set_opacity(window, 0.97)
add_text(
    s, Inches(6.68), Inches(1.42), Inches(5.15), Inches(0.36),
    [("CUSTOMIZABLE", 13, True, BLUE)],
)
add_text(
    s, Inches(6.68), Inches(1.82), Inches(5.15), Inches(0.72),
    [("Keep the controller.\nChange the vocabulary.", 23, True, WHITE)],
    space_after=1,
)
box(
    s, Inches(6.66), Inches(2.68), Inches(5.2), Inches(2.88),
    WHITE, radius=True,
)
add_picture(
    s, ACTION_MAPPING, Inches(6.78), Inches(2.79),
    Inches(4.96), Inches(2.62),
)
add_text(
    s, Inches(6.7), Inches(5.7), Inches(5.1), Inches(0.3),
    [("A-Z + 0-9  |  TESTED ACTIONS  |  SIM + REAL", 9, True, MIST)],
    align=PP_ALIGN.CENTER,
)
add_click_transition(s, effect="wipe", direction="d")
add_notes(s, """
[2:29-2:44]
Then we tailor the control language. Letters and digits can map to tested
actions. For each user or company, the reliable controller stays while the
task vocabulary adapts to their workflow.
""", cue="REMOTE: click to close the pitch.")


# 8. Open-ended closing and acknowledgements
s = make_slide(NAVY)
kicker(s, "Where this could go", dark=True)
add_text(
    s, Inches(0.78), Inches(1.04), Inches(11.8), Inches(1.0),
    [("One user. One task. One proof.", 40, True, WHITE)],
)
add_text(
    s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.48),
    [("Then repeat what genuinely helps.", 21, False, MIST)],
)
roadmap = [
    ("01", "LISTEN", "Find the real barrier."),
    ("02", "BUILD", "Co-design one useful task."),
    ("03", "PROVE", "Measure whether it helps."),
]
for index, (number, title, detail) in enumerate(roadmap):
    left = Inches(0.82 + index * 4.18)
    box(s, left, Inches(2.72), Inches(3.56), Inches(0.08), BLUE)
    add_text(
        s, left, Inches(3.06), Inches(0.58), Inches(0.42),
        [(number, 20, True, BLUE)],
    )
    add_text(
        s, left + Inches(0.66), Inches(3.08), Inches(2.78), Inches(0.34),
        [(title, 14, True, WHITE)],
    )
    add_text(
        s, left, Inches(3.72), Inches(3.5), Inches(0.68),
        [(detail, 17, False, MIST)],
        line_spacing=1.05,
    )
box(s, Inches(0.82), Inches(4.72), Inches(11.72), Inches(0.88), BLUE_DK,
    radius=True)
add_text(
    s, Inches(0.82), Inches(4.72), Inches(11.72), Inches(0.88),
    [("STARTUP PLAN v0.1:  HELP FIRST.  SCALE SECOND.", 17, True, WHITE)],
    align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE,
)
box(s, 0, Inches(6.12), SW, Inches(1.38), WHITE)
add_picture(
    s, TUM_LOGO, Inches(0.78), Inches(6.58), Inches(0.78), Inches(0.36)
)
add_picture(
    s, MIRMI_LOGO, Inches(1.82), Inches(6.58), Inches(1.1), Inches(0.36)
)
add_text(
    s, Inches(3.28), Inches(6.35), Inches(8.9), Inches(0.36),
    [("THANK YOU", 11, True, BLUE)],
)
add_text(
    s, Inches(3.28), Inches(6.7), Inches(8.9), Inches(0.36),
    [("Federico Masiero | Emanuele Aimi | Luca Borriello | "
      "Prof. Dr. Lorenzo Masia", 12, True, INK)],
)
add_text(
    s, Inches(11.9), Inches(7.08), Inches(0.7), Inches(0.25),
    [("8", 9, True, SUBTLE)],
    align=PP_ALIGN.RIGHT,
)
add_click_transition(s, effect="fade")
add_notes(s, """
[2:44-3:00]
MagPilot is a working platform, not a finished product. The next step is one
real user, one useful task, and one measurable benefit. If it works, we repeat.
Help first, scale second. Thank you.
""")


output = os.path.join(HERE, "MagPilot_Keynote.pptx")
prs.save(output)
print("Wrote", output, "-", len(prs.slides), "slides")
